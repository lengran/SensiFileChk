import os
import pytest

from src.checker import (
    _match_keywords,
    _extract_context,
    scan_single_file,
    scan_directory,
    Match,
    SUPPORTED_EXTENSIONS,
)


class TestMatchKeywords:
    def test_single_match(self):
        matches = _match_keywords("这是国家机密文件", ["国家机密"])
        assert len(matches) == 1
        assert matches[0].keyword == "国家机密"
        assert matches[0].start == 2

    def test_multiple_keywords(self):
        matches = _match_keywords("这是国家机密信息，包含内部资料", ["国家机密", "内部资料"])
        assert len(matches) == 2

    def test_overlapping_keywords(self):
        matches = _match_keywords("国家机密信息", ["国家机密", "机密信息"])
        assert len(matches) == 2

    def test_no_match(self):
        matches = _match_keywords("普通文本", ["机密"])
        assert len(matches) == 0

    def test_empty_keyword(self):
        matches = _match_keywords("文本", [""])
        assert len(matches) == 0

    def test_multiple_occurrences(self):
        matches = _match_keywords("机密文件和机密数据", ["机密"])
        assert len(matches) == 2

    def test_case_insensitive_english(self):
        """大小写不敏感匹配（英文关键词）"""
        matches = _match_keywords("This is SECRET information", ["secret"])
        assert len(matches) == 1
        assert matches[0].keyword == "secret"

    def test_case_insensitive_mixed(self):
        matches = _match_keywords("Secret and SECRET and secret", ["secret"])
        assert len(matches) == 3

    def test_chinese_keyword_still_case_sensitive(self):
        """中文关键词保持大小写敏感（中文无大小写）"""
        matches = _match_keywords("这是国家机密", ["国家机密"])
        assert len(matches) == 1

    def test_line_number_tracking(self):
        """行号追踪"""
        text = "第一行\n第二行机密信息\n第三行"
        matches = _match_keywords(text, ["机密信息"])
        assert len(matches) == 1
        assert matches[0].line_number == 2

    def test_line_number_first_line(self):
        text = "机密在开头"
        matches = _match_keywords(text, ["机密"])
        assert matches[0].line_number == 1


class TestExtractContext:
    def test_middle_context(self):
        text = "0123456789机密0123456789"
        ctx = _extract_context(text, 10, 12, 5)
        assert ctx == "56789机密01234"

    def test_start_boundary(self):
        text = "机密在开头"
        ctx = _extract_context(text, 0, 2, 5)
        assert ctx == "机密在开头"

    def test_end_boundary(self):
        text = "在结尾机密"
        ctx = _extract_context(text, 3, 5, 5)
        assert ctx == "在结尾机密"


class TestPhase1Discovery:
    def test_finds_supported_formats(self, tmp_path):
        (tmp_path / "a.txt").write_text("国家机密", encoding="utf-8")
        from docx import Document
        doc = Document()
        doc.add_paragraph("国家机密")
        doc.save(str(tmp_path / "b.docx"))
        result = scan_directory(str(tmp_path), ["国家机密"])
        assert len(result["results"]) + len(result["failures"]) == 2

    def test_ignores_unsupported(self, tmp_path):
        (tmp_path / "a.txt").write_text("国家机密", encoding="utf-8")
        (tmp_path / "b.jpg").write_bytes(b"\xff\xd8")
        result = scan_directory(str(tmp_path), ["国家机密"])
        total = len(result["results"]) + len(result["failures"])
        assert total == 1

    def test_recursive(self, tmp_path):
        sub = tmp_path / "sub"
        sub.mkdir()
        (tmp_path / "a.txt").write_text("国家机密", encoding="utf-8")
        (sub / "b.txt").write_text("国家机密", encoding="utf-8")
        result = scan_directory(str(tmp_path), ["国家机密"])
        assert len(result["results"]) == 2

    def test_empty_dir(self, tmp_path):
        result = scan_directory(str(tmp_path), ["机密"])
        assert result["results"] == []
        assert result["failures"] == []


class TestScanSingleFile:
    def test_scan_txt_with_match(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("这是国家机密文件", encoding="utf-8")
        result = scan_single_file(str(f), ["国家机密"], 50)
        assert len(result.matches) == 1
        assert result.error is None

    def test_scan_txt_no_match(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("普通文本", encoding="utf-8")
        result = scan_single_file(str(f), ["机密"], 50)
        assert len(result.matches) == 0

    def test_scan_unsupported_format(self, tmp_path):
        f = tmp_path / "test.xyz"
        f.write_text("内容")
        result = scan_single_file(str(f), ["内容"], 50)
        assert result.error is not None

    def test_scan_pdf(self, tmp_path):
        import fitz
        path = str(tmp_path / "test.pdf")
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "secret keyword in PDF", fontsize=12)
        doc.save(path)
        doc.close()
        result = scan_single_file(path, ["secret keyword"], 50, ocr_enabled=False)
        assert len(result.matches) >= 1

    def test_scan_docx(self, tmp_path):
        from docx import Document
        path = str(tmp_path / "test.docx")
        doc = Document()
        doc.add_paragraph("docx中的国家机密")
        doc.save(path)
        result = scan_single_file(path, ["国家机密"], 50)
        assert len(result.matches) == 1

    def test_scan_pptx(self, tmp_path):
        from pptx import Presentation
        path = str(tmp_path / "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "pptx中的国家机密"
        prs.save(path)
        result = scan_single_file(path, ["国家机密"], 50)
        assert len(result.matches) == 1

    def test_scan_xlsx(self, tmp_path):
        from openpyxl import Workbook
        path = str(tmp_path / "test.xlsx")
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "xlsx中的国家机密"
        wb.save(path)
        result = scan_single_file(path, ["国家机密"], 50)
        assert len(result.matches) == 1

    def test_case_insensitive_scan(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("This is SECRET data", encoding="utf-8")
        result = scan_single_file(str(f), ["secret"], 50)
        assert len(result.matches) == 1

    def test_scan_docx_with_table(self, tmp_path):
        from docx import Document
        path = str(tmp_path / "test_table.docx")
        doc = Document()
        doc.add_paragraph("普通段落")
        t = doc.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "表格中的国家机密"
        t.cell(0, 1).text = "普通内容"
        doc.save(path)
        result = scan_single_file(path, ["国家机密"], 50)
        assert len(result.matches) >= 1

    def test_scan_pptx_with_table(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        path = str(tmp_path / "test_table.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table_shape.table.cell(0, 0).text = "表格中的国家机密"
        table_shape.table.cell(0, 1).text = "普通内容"
        prs.save(path)
        result = scan_single_file(path, ["国家机密"], 50)
        assert len(result.matches) >= 1


class TestScanDirectory:
    def test_scan_with_matches(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("包含国家机密信息", encoding="utf-8")
        result = scan_directory(str(tmp_path), ["国家机密"])
        assert len(result["results"]) == 1
        assert len(result["failures"]) == 0

    def test_scan_no_keywords_match(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("普通文本", encoding="utf-8")
        result = scan_directory(str(tmp_path), ["机密"])
        assert len(result["results"]) == 0

    def test_scan_empty_dir(self, tmp_path):
        result = scan_directory(str(tmp_path), ["机密"])
        assert result["results"] == []
        assert result["failures"] == []

    def test_scan_mixed_formats(self, tmp_path):
        (tmp_path / "a.txt").write_text("国家机密", encoding="utf-8")
        from docx import Document
        doc = Document()
        doc.add_paragraph("内部资料")
        doc.save(str(tmp_path / "b.docx"))
        result = scan_directory(str(tmp_path), ["国家机密", "内部资料"])
        assert len(result["results"]) == 2

    def test_scan_with_failures(self, tmp_path):
        (tmp_path / "a.txt").write_text("国家机密", encoding="utf-8")
        corrupt = tmp_path / "corrupt.docx"
        corrupt.write_bytes(b"not a real docx")
        result = scan_directory(str(tmp_path), ["国家机密"])
        assert len(result["results"]) == 1
        assert len(result["failures"]) == 1


class TestParallelScan:
    def test_single_worker(self, tmp_path):
        for i in range(5):
            (tmp_path / f"f{i}.txt").write_text(f"文件{i}国家机密", encoding="utf-8")
        result = scan_directory(str(tmp_path), ["国家机密"], num_workers=1)
        assert len(result["results"]) == 5

    def test_multi_worker(self, tmp_path):
        for i in range(10):
            (tmp_path / f"f{i}.txt").write_text(f"文件{i}国家机密", encoding="utf-8")
        result = scan_directory(str(tmp_path), ["国家机密"], num_workers=2)
        assert len(result["results"]) == 10


class TestVerboseMode:
    def test_verbose_progress(self, tmp_path, capsys):
        for i in range(5):
            (tmp_path / f"f{i}.txt").write_text(f"文件{i}国家机密", encoding="utf-8")
        scan_directory(str(tmp_path), ["国家机密"], verbose=True)
        out = capsys.readouterr().out
        assert "扫描中" in out

    def test_verbose_hits(self, tmp_path, capsys):
        for i in range(5):
            (tmp_path / f"f{i}.txt").write_text(f"文件{i}国家机密", encoding="utf-8")
        scan_directory(str(tmp_path), ["国家机密"], verbose=True)
        out = capsys.readouterr().out
        assert out.count("[命中]") == 5

    def test_verbose_failures(self, tmp_path, capsys):
        (tmp_path / "f1.txt").write_text("正常", encoding="utf-8")
        corrupt = tmp_path / "f2.docx"
        corrupt.write_bytes(b"not a docx")
        scan_directory(str(tmp_path), ["正常"], verbose=True)
        out = capsys.readouterr().out
        assert "失败" in out

    def test_verbose_multi_worker(self, tmp_path, capsys):
        for i in range(5):
            (tmp_path / f"f{i}.txt").write_text(f"文件{i}国家机密", encoding="utf-8")
        scan_directory(str(tmp_path), ["国家机密"], num_workers=2, verbose=True)
        out = capsys.readouterr().out
        assert "命中" in out


class TestScanSingleFileAdvanced:
    def test_scan_corrupted_pdf(self, tmp_path):
        f = tmp_path / "bad.pdf"
        f.write_bytes(b"not a pdf")
        result = scan_single_file(str(f), ["test"], 50)
        assert result.error is not None

    def test_scan_parser_exception(self, tmp_path, monkeypatch):
        f = tmp_path / "test.txt"
        f.write_text("内容", encoding="utf-8")
        from unittest.mock import MagicMock
        monkeypatch.setattr("src.checker.TxtParser", lambda: MagicMock(parse=MagicMock(side_effect=Exception("boom"))))
        result = scan_single_file(str(f), ["内容"], 50)
        assert "解析失败" in result.error

    def test_scan_check_archives_false(self, tmp_path):
        import zipfile
        inner = tmp_path / "inner.txt"
        inner.write_text("机密", encoding="utf-8")
        zip_path = str(tmp_path / "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(str(inner), "inner.txt")
        result = scan_single_file(zip_path, ["机密"], 50, check_archives=False)
        assert result.error is not None

    def test_scan_zip_with_match(self, tmp_path):
        import zipfile
        inner = tmp_path / "inner.txt"
        inner.write_text("国家机密在压缩包内", encoding="utf-8")
        zip_path = str(tmp_path / "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(str(inner), "inner.txt")
        result = scan_single_file(zip_path, ["国家机密"], 50)
        assert len(result.matches) >= 1


class TestMultiWorkerEdgeCases:
    def test_multi_worker_results_match(self, tmp_path):
        for i in range(10):
            (tmp_path / f"f{i}.txt").write_text(f"文件{i}国家机密", encoding="utf-8")
        r1 = scan_directory(str(tmp_path), ["国家机密"], num_workers=1)
        r2 = scan_directory(str(tmp_path), ["国家机密"], num_workers=2)
        paths1 = sorted(fr.file_path for fr in r1["results"])
        paths2 = sorted(fr.file_path for fr in r2["results"])
        assert paths1 == paths2

    def test_multi_worker_failure(self, tmp_path):
        corrupt = tmp_path / "bad.docx"
        corrupt.write_bytes(b"not a docx")
        result = scan_directory(str(tmp_path), ["test"], num_workers=2)
        assert len(result["failures"]) == 1


class TestParserByExt:
    def test_pdf_parser(self):
        from src.checker import _parser_by_ext
        parser = _parser_by_ext(".pdf", ocr_enabled=True)
        assert parser is not None
        assert parser.ocr_enabled is True

    def test_office_parser(self):
        from src.checker import _parser_by_ext
        parser = _parser_by_ext(".docx")
        assert parser is not None

    def test_txt_parser(self):
        from src.checker import _parser_by_ext
        parser = _parser_by_ext(".txt")
        assert parser is not None

    def test_unknown_ext(self):
        from src.checker import _parser_by_ext
        parser = _parser_by_ext(".xyz")
        assert parser is None


class TestWorkerTask:
    def test_worker_task_delegates(self, tmp_path):
        from src.checker import _worker_task
        f = tmp_path / "test.txt"
        f.write_text("国家机密", encoding="utf-8")
        result = _worker_task((str(f), ["国家机密"], 50, False, True))
        assert len(result.matches) == 1


class TestMultiWorkerVerbose:
    def test_verbose_multi_worker_failure(self, tmp_path, capsys):
        corrupt = tmp_path / "bad.docx"
        corrupt.write_bytes(b"not a docx")
        (tmp_path / "good.txt").write_text("正常文本", encoding="utf-8")
        scan_directory(str(tmp_path), ["正常文本"], verbose=True, num_workers=2)
        out = capsys.readouterr().out
        assert "失败" in out

    def test_worker_exception_in_pool(self, tmp_path, monkeypatch):
        from unittest.mock import MagicMock, patch
        (tmp_path / "f.txt").write_text("国家机密", encoding="utf-8")
        with patch("src.checker._worker_task", side_effect=RuntimeError("worker crashed")):
            result = scan_directory(str(tmp_path), ["国家机密"], num_workers=2)
            assert len(result["failures"]) == 1
            assert "worker 异常" in result["failures"][0].error


class TestTwoPhaseProgress:
    def test_discovery_phase_output(self, tmp_path, capsys):
        for i in range(3):
            (tmp_path / f"f{i}.txt").write_text("正常", encoding="utf-8")
        scan_directory(str(tmp_path), ["机密"], verbose=False)
        out = capsys.readouterr().out
        assert "发现文件中" in out

    def test_scan_phase_with_percentage(self, tmp_path, capsys):
        for i in range(3):
            (tmp_path / f"f{i}.txt").write_text("国家机密", encoding="utf-8")
        scan_directory(str(tmp_path), ["国家机密"], verbose=False)
        out = capsys.readouterr().out
        assert "扫描中" in out
        assert "已检测: 3/3" in out
        assert "100.0%" in out

    def test_discovery_phase_multi_worker(self, tmp_path, capsys):
        for i in range(3):
            (tmp_path / f"f{i}.txt").write_text("国家机密", encoding="utf-8")
        scan_directory(str(tmp_path), ["国家机密"], num_workers=2, verbose=False)
        out = capsys.readouterr().out
        assert "发现文件中" in out
        assert "扫描中" in out


class TestPhase1ZeroStat:
    def test_single_worker_no_stat_calls(self, tmp_path, monkeypatch):
        for i in range(3):
            (tmp_path / f"f{i}.txt").write_text(f"国家机密{i}", encoding="utf-8")
        call_count = 0
        original_getsize = os.path.getsize

        def counting_getsize(p):
            nonlocal call_count
            call_count += 1
            return original_getsize(p)

        monkeypatch.setattr(os.path, "getsize", counting_getsize)
        result = scan_directory(str(tmp_path), ["国家机密"], num_workers=1)
        assert len(result["results"]) == 3
        assert call_count == 0

    def test_multi_worker_stat_only_in_phase2(self, tmp_path, monkeypatch):
        for i in range(5):
            (tmp_path / f"f{i}.txt").write_text(f"国家机密{i}", encoding="utf-8")
        original_getsize = os.path.getsize
        phase1_calls = 0
        phase2_calls = 0
        in_phase2 = False

        def tracking_getsize(p):
            nonlocal phase1_calls, phase2_calls, in_phase2
            if in_phase2:
                phase2_calls += 1
            else:
                phase1_calls += 1
            return original_getsize(p)

        monkeypatch.setattr(os.path, "getsize", tracking_getsize)

        from unittest.mock import patch
        original_walk = os.walk

        def tracking_walk(d):
            nonlocal in_phase2
            yield from original_walk(d)
            in_phase2 = True

        with patch("src.checker.os.walk", tracking_walk):
            result = scan_directory(str(tmp_path), ["国家机密"], num_workers=2)
        assert len(result["results"]) == 5
        assert phase1_calls == 0
        assert phase2_calls == 5


class TestLargeFileInline:
    def test_large_file_goes_inline(self, tmp_path, monkeypatch):
        (tmp_path / "small.txt").write_text("国家机密", encoding="utf-8")
        (tmp_path / "big.txt").write_text("国家机密在大文件中", encoding="utf-8")
        monkeypatch.setattr(os.path, "getsize", lambda p: 600 * 1024 * 1024 if "big" in p else 10)
        result = scan_directory(str(tmp_path), ["国家机密"], num_workers=2)
        assert len(result["results"]) == 2

    def test_large_file_single_worker(self, tmp_path, monkeypatch):
        (tmp_path / "big.txt").write_text("国家机密", encoding="utf-8")
        monkeypatch.setattr(os.path, "getsize", lambda p: 600 * 1024 * 1024)
        result = scan_directory(str(tmp_path), ["国家机密"], num_workers=1)
        assert len(result["results"]) == 1


class TestGetFileSize:
    def test_regular_file(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("hello", encoding="utf-8")
        from src.checker import _get_file_size
        assert _get_file_size(str(f)) == 5

    def test_nonexistent_file(self):
        from src.checker import _get_file_size
        assert _get_file_size("/nonexistent/file.txt") == 0


class TestEstimateBytesFromSize:
    def test_regular_file(self):
        from src.checker import _estimate_bytes_from_size
        assert _estimate_bytes_from_size(100, ".txt") == 100

    def test_archive_file_multiplier(self):
        from src.checker import _estimate_bytes_from_size
        assert _estimate_bytes_from_size(100, ".zip") == 500

    def test_targz_multiplier(self):
        from src.checker import _estimate_bytes_from_size
        assert _estimate_bytes_from_size(200, ".tgz") == 1000

    def test_zero_size(self):
        from src.checker import _estimate_bytes_from_size
        assert _estimate_bytes_from_size(0, ".txt") == 0
        assert _estimate_bytes_from_size(0, ".zip") == 0


class TestMemoryThrottling:
    def test_throttle_with_low_budget(self, tmp_path, monkeypatch):
        from src.checker import MAX_CONCURRENT_BYTES
        for i in range(5):
            (tmp_path / f"f{i}.txt").write_text(f"国家机密{i}", encoding="utf-8")
        monkeypatch.setattr("src.checker.MAX_CONCURRENT_BYTES", 1)
        result = scan_directory(str(tmp_path), ["国家机密"], num_workers=2)
        assert len(result["results"]) == 5


class TestBrokenPoolRecovery:
    def test_broken_pool_moves_to_inline(self, tmp_path, monkeypatch):
        from unittest.mock import patch, MagicMock
        from concurrent.futures import ProcessPoolExecutor as _PPE
        from concurrent.futures.process import BrokenProcessPool
        for i in range(3):
            (tmp_path / f"f{i}.txt").write_text(f"国家机密{i}", encoding="utf-8")

        original_submit = _PPE.submit

        call_count = 0

        def broken_submit(self, fn, *args, **kwargs):
            nonlocal call_count
            call_count += 1
            if call_count <= 2:
                f = original_submit(self, fn, *args, **kwargs)
                return f
            raise BrokenProcessPool("test broken pool")

        with patch.object(_PPE, "submit", broken_submit):
            result = scan_directory(str(tmp_path), ["国家机密"], num_workers=2)

        total = len(result["results"]) + len(result["failures"])
        assert total == 3
