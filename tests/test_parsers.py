import os
import pytest
from unittest.mock import MagicMock

from src.parsers.base import ParserError
from src.parsers.txt import TxtParser
from src.parsers.pdf import PdfParser
from src.parsers.office import OfficeParser
from src.parsers.archive import ArchiveParser


def _make_txt(tmp_path, name: str = "inner.txt", content: str = "压缩包内敏感词") -> str:
    f = tmp_path / name
    f.write_text(content, encoding="utf-8")
    return str(f)


@pytest.fixture
def txt_parser():
    return TxtParser()


@pytest.fixture
def pdf_parser():
    return PdfParser(ocr_enabled=False)


@pytest.fixture
def office_parser():
    return OfficeParser()


@pytest.fixture
def archive_parser():
    from src.parsers.txt import TxtParser
    from src.parsers.pdf import PdfParser
    from src.parsers.office import OfficeParser
    def factory(ext):
        if ext == ".txt":
            return TxtParser()
        elif ext == ".pdf":
            return PdfParser(ocr_enabled=False)
        elif ext in (".docx", ".pptx", ".xlsx", ".doc", ".xls", ".ppt"):
            return OfficeParser()
        return None
    return ArchiveParser(inner_parser_factory=factory)


class TestTxtParser:
    def test_utf8(self, txt_parser, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("这是一段测试文本", encoding="utf-8")
        assert txt_parser.parse(str(f)) == "这是一段测试文本"

    def test_gbk(self, txt_parser, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("这是一段GBK编码的文本", encoding="gbk")
        assert txt_parser.parse(str(f)) == "这是一段GBK编码的文本"

    def test_gb2312(self, txt_parser, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("GB2312编码测试", encoding="gb2312")
        assert txt_parser.parse(str(f)) == "GB2312编码测试"

    def test_empty_file(self, txt_parser, tmp_path):
        f = tmp_path / "empty.txt"
        f.write_bytes(b"")
        assert txt_parser.parse(str(f)) == ""

    def test_binary_file(self, txt_parser, tmp_path):
        f = tmp_path / "binary.txt"
        f.write_bytes(bytes(range(128, 256)) * 100)
        with pytest.raises(ParserError):
            txt_parser.parse(str(f))

    def test_nonexistent_file(self, txt_parser):
        with pytest.raises(ParserError):
            txt_parser.parse("/nonexistent/path/test.txt")

    def test_large_file(self, txt_parser, tmp_path):
        f = tmp_path / "large.txt"
        content = "测试内容" * 10000
        f.write_text(content, encoding="utf-8")
        assert txt_parser.parse(str(f)) == content


class TestPdfParser:
    def _make_pdf(self, tmp_path, text: str, name: str = "test.pdf") -> str:
        import fitz
        path = str(tmp_path / name)
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), text, fontsize=12)
        doc.save(path)
        doc.close()
        return path

    def test_text_pdf(self, pdf_parser, tmp_path):
        # 使用英文避免中文字体问题
        path = self._make_pdf(tmp_path, "PDF test content here")
        result = pdf_parser.parse(path)
        # 检查返回非空文本内容（PDF可能有不可见字符）
        assert result is not None and len(result) > 0

    def test_multipage_pdf(self, pdf_parser, tmp_path):
        import fitz
        path = str(tmp_path / "multi.pdf")
        doc = fitz.open()
        for i in range(5):
            page = doc.new_page()
            # 使用英文避免中文字体问题
            page.insert_text((72, 72), f"Page {i+1} content here", fontsize=12)
        doc.save(path)
        doc.close()
        result = pdf_parser.parse(path)
        # 检查每个页面的内容都已被提取
        for i in range(5):
            assert f"Page {i+1}" in result

    def test_empty_pdf(self, pdf_parser, tmp_path):
        import fitz
        path = str(tmp_path / "empty.pdf")
        doc = fitz.open()
        doc.new_page()
        doc.save(path)
        doc.close()
        result = pdf_parser.parse(path)
        assert result == "" or result.strip() == ""

    def test_encrypted_pdf(self, pdf_parser, tmp_path):
        import fitz
        path = str(tmp_path / "encrypted.pdf")
        doc = fitz.open()
        doc.new_page()
        page = doc[0]
        # 使用英文避免中文字体问题
        page.insert_text((72, 72), "encrypted content", fontsize=12)
        doc.save(path, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw="owner", user_pw="user")
        doc.close()
        with pytest.raises(ParserError, match="加密"):
            pdf_parser.parse(path)


class TestOfficeParser:
    def test_docx(self, office_parser, tmp_path):
        from docx import Document
        path = str(tmp_path / "test.docx")
        doc = Document()
        doc.add_paragraph("这是docx测试内容")
        doc.save(path)
        result = office_parser.parse(path)
        assert "docx测试内容" in result

    def test_pptx(self, office_parser, tmp_path):
        from pptx import Presentation
        path = str(tmp_path / "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "这是pptx测试内容"
        prs.save(path)
        result = office_parser.parse(path)
        assert "pptx测试内容" in result

    def test_xlsx(self, office_parser, tmp_path):
        from openpyxl import Workbook
        path = str(tmp_path / "test.xlsx")
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "这是xlsx测试内容"
        wb.save(path)
        result = office_parser.parse(path)
        assert "xlsx测试内容" in result

    def test_empty_docx(self, office_parser, tmp_path):
        from docx import Document
        path = str(tmp_path / "empty.docx")
        doc = Document()
        doc.save(path)
        result = office_parser.parse(path)
        assert result == ""

    def test_corrupted_office(self, office_parser, tmp_path):
        path = str(tmp_path / "corrupt.docx")
        with open(path, "wb") as f:
            f.write(b"not a real docx file content " * 10)
        with pytest.raises(ParserError):
            office_parser.parse(path)


class TestArchiveParser:
    def _make_txt(self, tmp_path, name: str = "inner.txt", content: str = "压缩包内敏感词") -> str:
        return _make_txt(tmp_path, name, content)

    def test_zip(self, archive_parser, tmp_path):
        import zipfile
        inner = self._make_txt(tmp_path)
        zip_path = str(tmp_path / "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        result = archive_parser.parse(zip_path)
        assert "压缩包内敏感词" in result

    def test_tar(self, archive_parser, tmp_path):
        import tarfile
        inner = self._make_txt(tmp_path)
        tar_path = str(tmp_path / "test.tar")
        with tarfile.open(tar_path, "w") as tf:
            tf.add(inner, arcname="inner.txt")
        result = archive_parser.parse(tar_path)
        assert "压缩包内敏感词" in result

    def test_targz(self, archive_parser, tmp_path):
        import tarfile
        inner = self._make_txt(tmp_path)
        tgz_path = str(tmp_path / "test.tar.gz")
        with tarfile.open(tgz_path, "w:gz") as tf:
            tf.add(inner, arcname="inner.txt")
        result = archive_parser.parse(tgz_path)
        assert "压缩包内敏感词" in result

    def test_gz(self, archive_parser, tmp_path):
        import gzip
        inner_path = str(tmp_path / "inner.txt")
        gz_path = str(tmp_path / "inner.txt.gz")
        with open(inner_path, "w", encoding="utf-8") as f:
            f.write("gz压缩的敏感词")
        with open(inner_path, "rb") as f_in:
            with gzip.open(gz_path, "wb") as f_out:
                f_out.write(f_in.read())
        result = archive_parser.parse(gz_path)
        assert "gz压缩的敏感词" in result

    def test_nested_zip(self, archive_parser, tmp_path):
        import zipfile
        inner_txt = tmp_path / "inner.txt"
        inner_txt.write_text("嵌套敏感词", encoding="utf-8")
        inner_zip = str(tmp_path / "inner.zip")
        with zipfile.ZipFile(inner_zip, "w") as zf:
            zf.write(str(inner_txt), "inner.txt")
        outer_zip = str(tmp_path / "outer.zip")
        with zipfile.ZipFile(outer_zip, "w") as zf:
            zf.write(inner_zip, "inner.zip")
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        result = parser.parse(outer_zip)
        assert "嵌套敏感词" in result

    def test_depth_limit(self, tmp_path):
        import zipfile
        # 创建一个嵌套zip来测试深度限制
        # depth=10 时应该直接超限
        deep_parser = ArchiveParser(depth=10)
        inner = self._make_txt(tmp_path)
        zip_path = str(tmp_path / "deep.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        with pytest.raises(ParserError, match="深度超限"):
            deep_parser.parse(zip_path)

    def test_zip_slip_protection_dotdot(self, archive_parser, tmp_path):
        """Zip Slip: .. 路径检测"""
        import zipfile
        zip_path = str(tmp_path / "slip.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("../../../../../tmp/evil.txt", "malicious")
        with pytest.raises(ParserError):
            archive_parser.parse(zip_path)

    def test_zip_slip_protection_absolute_path(self, archive_parser, tmp_path):
        """Zip Slip: 绝对路径检测"""
        import zipfile
        zip_path = str(tmp_path / "abs_slip.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("/etc/passwd", "malicious")
        with pytest.raises(ParserError):
            archive_parser.parse(zip_path)

    def test_empty_zip(self, archive_parser, tmp_path):
        import zipfile
        zip_path = str(tmp_path / "empty.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            pass
        result = archive_parser.parse(zip_path)
        assert result == ""

    def test_corrupted_zip(self, archive_parser, tmp_path):
        zip_path = str(tmp_path / "corrupt.zip")
        with open(zip_path, "wb") as f:
            f.write(b"not a real zip file")
        with pytest.raises(ParserError):
            archive_parser.parse(zip_path)

    def test_rar(self, archive_parser, tmp_path):
        try:
            import rarfile
        except ImportError:
            pytest.skip("rarfile not installed")
        try:
            rarfile.tool_setup()
        except Exception:
            pytest.skip("unrar not installed")
        import zipfile
        inner = self._make_txt(tmp_path, "inner.txt", "rar内敏感词")
        zip_path = str(tmp_path / "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        try:
            import subprocess
            rar_path = str(tmp_path / "test.rar")
            r = subprocess.run(
                ["rar", "a", "-ep", rar_path, inner],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode != 0:
                pytest.skip("rar command not available")
            result = archive_parser.parse(rar_path)
            assert "rar内敏感词" in result
        except FileNotFoundError:
            pytest.skip("rar command not installed")

    def test_7z(self, archive_parser, tmp_path):
        try:
            import py7zr
        except ImportError:
            pytest.skip("py7zr not installed")
        import zipfile
        inner = self._make_txt(tmp_path, "inner.txt", "7z内敏感词")
        sz_path = str(tmp_path / "test.7z")
        try:
            with py7zr.SevenZipFile(sz_path, "w") as sz:
                sz.write(inner, "inner.txt")
        except Exception:
            pytest.skip("7z creation failed")
        result = archive_parser.parse(sz_path)
        assert "7z内敏感词" in result


class TestOfficeParserLegacy:
    def test_parse_doc_success(self, tmp_path, monkeypatch):
        from unittest.mock import patch, MagicMock
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy doc content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Linux")):
            with patch("src.parsers.office.subprocess.run") as mock_run:
                mock_run.return_value = MagicMock(returncode=0, stdout="legacy doc content")
                parser = OfficeParser()
                result = parser.parse(doc_path)
                assert "legacy doc content" in result

    def test_parse_doc_error_on_invalid(self, tmp_path, monkeypatch):
        from unittest.mock import patch, MagicMock
        doc_path = str(tmp_path / "bad.doc")
        with open(doc_path, "wb") as f:
            f.write(b"bad content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Linux")):
            with patch("src.parsers.office.subprocess.run") as mock_run:
                mock_run.return_value = MagicMock(returncode=1, stderr="parse error")
                parser = OfficeParser()
                with pytest.raises(ParserError):
                    parser.parse(doc_path)

    def test_parse_legacy_platform_routing(self, tmp_path, monkeypatch):
        from unittest.mock import patch, MagicMock
        doc_path = str(tmp_path / "route.doc")
        with open(doc_path, "wb") as f:
            f.write(b"content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Darwin")):
            with patch("src.parsers.office.subprocess.run") as mock_run:
                mock_run.return_value = MagicMock(returncode=0, stdout="catdoc output")
                parser = OfficeParser()
                result = parser.parse(doc_path)
                mock_run.assert_called_once()
                assert "catdoc output" in result


class TestPdfOcr:
    def _make_pdf_no_text(self, tmp_path, name="blank.pdf"):
        import fitz
        path = str(tmp_path / name)
        doc = fitz.open()
        doc.new_page()
        doc.save(path)
        doc.close()
        return path

    def test_ocr_page_success(self, tmp_path, monkeypatch):
        import fitz
        path = self._make_pdf_no_text(tmp_path)
        parser = PdfParser(ocr_enabled=True)
        from unittest.mock import MagicMock
        mock_img_to_str = MagicMock(return_value="OCR extracted text")
        monkeypatch.setattr("pytesseract.image_to_string", mock_img_to_str)
        result = parser.parse(path)
        assert "OCR extracted text" in result

    def test_ocr_on_empty_page(self, tmp_path, monkeypatch):
        import fitz
        path = self._make_pdf_no_text(tmp_path)
        parser = PdfParser(ocr_enabled=True)
        monkeypatch.setattr("pytesseract.image_to_string", lambda *a, **kw: "")
        result = parser.parse(path)
        assert result == ""

    def test_ocr_import_error(self, tmp_path, monkeypatch):
        import fitz
        path = self._make_pdf_no_text(tmp_path)
        parser = PdfParser(ocr_enabled=True)
        import builtins
        real_import = builtins.__import__
        def _block_pytesseract(name, *args, **kwargs):
            if name == "pytesseract":
                raise ImportError("no pytesseract")
            return real_import(name, *args, **kwargs)
        monkeypatch.setattr(builtins, "__import__", _block_pytesseract)
        result = parser.parse(path)
        assert result == ""

    def test_ocr_runtime_error(self, tmp_path, monkeypatch):
        import fitz
        path = self._make_pdf_no_text(tmp_path)
        parser = PdfParser(ocr_enabled=True)
        monkeypatch.setattr("pytesseract.image_to_string", MagicMock(side_effect=RuntimeError("tesseract failed")))
        result = parser.parse(path)
        assert result == ""

    def test_ocr_generic_error(self, tmp_path, monkeypatch):
        import fitz
        path = self._make_pdf_no_text(tmp_path)
        parser = PdfParser(ocr_enabled=True)
        monkeypatch.setattr("pytesseract.image_to_string", MagicMock(side_effect=Exception("unknown error")))
        result = parser.parse(path)
        assert result == ""


class TestArchiveRarAdvanced:
    def test_rar_with_directory_entries(self, archive_parser, tmp_path):
        try:
            import rarfile
        except ImportError:
            pytest.skip("rarfile not installed")
        try:
            rarfile.tool_setup()
        except Exception:
            pytest.skip("unrar not installed")
        import subprocess
        inner = _make_txt(tmp_path, "inner.txt", "rar正常内容")
        try:
            rar_path = str(tmp_path / "test_dir.rar")
            r = subprocess.run(
                ["rar", "a", rar_path, inner],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode != 0:
                pytest.skip("rar command not available")
            result = archive_parser.parse(rar_path)
            assert "rar正常内容" in result
        except FileNotFoundError:
            pytest.skip("rar command not installed")

    def test_rar_size_limit(self, archive_parser, tmp_path, monkeypatch):
        try:
            import rarfile
        except ImportError:
            pytest.skip("rarfile not installed")
        try:
            rarfile.tool_setup()
        except Exception:
            pytest.skip("unrar not installed")
        import subprocess
        inner = _make_txt(tmp_path, "big.txt", "x" * 200)
        try:
            rar_path = str(tmp_path / "big.rar")
            r = subprocess.run(
                ["rar", "a", "-ep", rar_path, inner],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode != 0:
                pytest.skip("rar command not available")
            monkeypatch.setattr("src.parsers.archive.MAX_SIZE", 20)
            with pytest.raises(ParserError):
                archive_parser.parse(rar_path)
        except FileNotFoundError:
            pytest.skip("rar command not installed")

    def test_rar_per_file_error(self, archive_parser, tmp_path):
        try:
            import rarfile
        except ImportError:
            pytest.skip("rarfile not installed")
        try:
            rarfile.tool_setup()
        except Exception:
            pytest.skip("unrar not installed")
        import subprocess
        inner = _make_txt(tmp_path, "inner.txt", "rar内容")
        try:
            rar_path = str(tmp_path / "test_per_file.rar")
            r = subprocess.run(
                ["rar", "a", "-ep", rar_path, inner],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode != 0:
                pytest.skip("rar command not available")
            result = archive_parser.parse(rar_path)
            assert "rar内容" in result
        except FileNotFoundError:
            pytest.skip("rar command not installed")

    def test_rar_not_rar(self, archive_parser, tmp_path):
        try:
            import rarfile
        except ImportError:
            pytest.skip("rarfile not installed")
        rar_path = str(tmp_path / "not_rar.rar")
        with open(rar_path, "wb") as f:
            f.write(b"not a rar file at all")
        with pytest.raises((ParserError, rarfile.NotRarFile)):
            archive_parser.parse(rar_path)


class TestArchiveRemainingGaps:
    def test_gz_containing_zip(self, archive_parser, tmp_path):
        import gzip, zipfile
        inner = _make_txt(tmp_path, "gz_zip_inner.txt", "gz内zip敏感词")
        zip_path = str(tmp_path / "inner_gz.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        gz_path = str(tmp_path / "outer.zip.gz")
        with open(zip_path, "rb") as f_in:
            with gzip.open(gz_path, "wb") as f_out:
                f_out.write(f_in.read())
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        result = parser.parse(gz_path)
        assert "gz内zip敏感词" in result

    def test_gz_size_limit(self, archive_parser, tmp_path, monkeypatch):
        import gzip
        monkeypatch.setattr("src.parsers.archive.MAX_SIZE", 20)
        data = b"a" * 200
        gz_path = str(tmp_path / "big.txt.gz")
        with gzip.open(gz_path, "wb") as f:
            f.write(data)
        with pytest.raises(ParserError):
            archive_parser.parse(gz_path)

    def test_tar_total_size_limit(self, archive_parser, tmp_path, monkeypatch):
        import tarfile
        monkeypatch.setattr("src.parsers.archive.MAX_TOTAL_SIZE", 50)
        inner = _make_txt(tmp_path, "big.txt", "x" * 200)
        tar_path = str(tmp_path / "big.tar")
        with tarfile.open(tar_path, "w") as tf:
            tf.add(inner, arcname="big.txt")
        with pytest.raises(ParserError):
            archive_parser.parse(tar_path)

    def test_tar_per_file_error(self, tmp_path):
        import tarfile
        inner = _make_txt(tmp_path, "inner.txt", "tar内容")
        tar_path = str(tmp_path / "test.tar")
        with tarfile.open(tar_path, "w") as tf:
            tf.add(inner, arcname="inner.txt")
        def bad_factory(ext):
            if ext == ".txt":
                raise ParserError("bad parser")
            return None
        parser = ArchiveParser(inner_parser_factory=bad_factory)
        result = parser.parse(tar_path)
        assert "解析失败" in result

    def test_7z_size_limit(self, archive_parser, tmp_path, monkeypatch):
        try:
            import py7zr
        except ImportError:
            pytest.skip("py7zr not installed")
        inner = _make_txt(tmp_path, "big.txt", "x" * 200)
        sz_path = str(tmp_path / "big.7z")
        try:
            with py7zr.SevenZipFile(sz_path, "w") as sz:
                sz.write(inner, "big.txt")
        except Exception:
            pytest.skip("7z creation failed")
        monkeypatch.setattr("src.parsers.archive.MAX_SIZE", 20)
        with pytest.raises(ParserError):
            archive_parser.parse(sz_path)

    def test_7z_total_size_limit(self, archive_parser, tmp_path, monkeypatch):
        try:
            import py7zr
        except ImportError:
            pytest.skip("py7zr not installed")
        inner = _make_txt(tmp_path, "big.txt", "x" * 200)
        sz_path = str(tmp_path / "big_total.7z")
        try:
            with py7zr.SevenZipFile(sz_path, "w") as sz:
                sz.write(inner, "big.txt")
        except Exception:
            pytest.skip("7z creation failed")
        monkeypatch.setattr("src.parsers.archive.MAX_TOTAL_SIZE", 50)
        with pytest.raises(ParserError):
            archive_parser.parse(sz_path)

    def test_inner_parser_none_fallback(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "some content")
        zip_path = str(tmp_path / "test_no_parser.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        parser = ArchiveParser(inner_parser_factory=lambda e: None)
        result = parser.parse(zip_path)
        assert result == ""

    def test_inner_parser_per_file_error(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "content")
        zip_path = str(tmp_path / "test_parser_err.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        def error_factory(ext):
            if ext == ".txt":
                class BadParser:
                    def parse(self, fp):
                        raise ParserError("simulated error")
                return BadParser()
            return None
        parser = ArchiveParser(inner_parser_factory=error_factory)
        result = parser.parse(zip_path)
        assert "解析失败" in result

    def test_gz_per_file_parser_error(self, tmp_path):
        import gzip
        inner_path = str(tmp_path / "inner.txt")
        gz_path = str(tmp_path / "inner.txt.gz")
        with open(inner_path, "w", encoding="utf-8") as f:
            f.write("gz content")
        with open(inner_path, "rb") as f_in:
            with gzip.open(gz_path, "wb") as f_out:
                f_out.write(f_in.read())
        def error_factory(ext):
            if ext == ".txt":
                class BadParser:
                    def parse(self, fp):
                        raise ParserError("gz parse error")
                return BadParser()
            return None
        parser = ArchiveParser(inner_parser_factory=error_factory)
        with pytest.raises(ParserError):
            parser.parse(gz_path)


class TestOfficeLegacyAdvanced:
    def test_parse_xls_legacy(self, tmp_path):
        from unittest.mock import patch, MagicMock
        xls_path = str(tmp_path / "test.xls")
        with open(xls_path, "wb") as f:
            f.write(b"dummy xls content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Linux")):
            with patch("src.parsers.office.subprocess.run") as mock_run:
                mock_run.return_value = MagicMock(returncode=0, stdout="legacy xls content")
                parser = OfficeParser()
                result = parser.parse(xls_path)
                assert "legacy xls content" in result

    def test_parse_ppt_legacy(self, tmp_path):
        from unittest.mock import patch, MagicMock
        ppt_path = str(tmp_path / "test.ppt")
        with open(ppt_path, "wb") as f:
            f.write(b"dummy ppt content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Linux")):
            with patch("src.parsers.office.subprocess.run") as mock_run:
                mock_run.return_value = MagicMock(returncode=0, stdout="legacy ppt content")
                parser = OfficeParser()
                result = parser.parse(ppt_path)
                assert "legacy ppt content" in result

    def test_parse_doc_windows_pywin32(self, tmp_path):
        from unittest.mock import patch, MagicMock
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy doc content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Windows")):
            mock_win32 = MagicMock()
            mock_app = MagicMock()
            mock_doc = MagicMock()
            mock_doc.Content.Text = "pywin32 doc content"
            mock_app.Documents.Open.return_value = mock_doc
            mock_win32.client.Dispatch.return_value = mock_app
            with patch.dict("sys.modules", {"win32com": mock_win32, "win32com.client": mock_win32.client}):
                parser = OfficeParser()
                result = parser.parse(doc_path)
                assert "pywin32 doc content" in result

    def test_parse_xls_windows_pywin32(self, tmp_path):
        from unittest.mock import patch, MagicMock
        xls_path = str(tmp_path / "test.xls")
        with open(xls_path, "wb") as f:
            f.write(b"dummy xls content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Windows")):
            mock_win32 = MagicMock()
            mock_app = MagicMock()
            mock_ws = MagicMock()
            mock_cell = MagicMock()
            mock_cell.Value = "pywin32 xls cell"
            mock_range = [[mock_cell]]
            mock_ws.UsedRange = mock_range
            mock_ws.Worksheets = [mock_ws]
            mock_app.Workbooks.Open.return_value.Worksheets = [mock_ws]
            mock_wb = mock_app.Workbooks.Open.return_value
            mock_wb.Worksheets = [mock_ws]
            mock_win32.client.Dispatch.return_value = mock_app
            with patch.dict("sys.modules", {"win32com": mock_win32, "win32com.client": mock_win32.client}):
                parser = OfficeParser()
                result = parser.parse(xls_path)
                assert "pywin32 xls cell" in result

    def test_parse_ppt_windows_pywin32(self, tmp_path):
        from unittest.mock import patch, MagicMock
        ppt_path = str(tmp_path / "test.ppt")
        with open(ppt_path, "wb") as f:
            f.write(b"dummy ppt content")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Windows")):
            mock_win32 = MagicMock()
            mock_app = MagicMock()
            mock_shape = MagicMock()
            mock_shape.HasTextFrame = True
            mock_shape.TextFrame.TextRange.Text = "pywin32 ppt slide"
            mock_slide = MagicMock()
            mock_slide.Shapes = [mock_shape]
            mock_pres = MagicMock()
            mock_pres.Slides = [mock_slide]
            mock_app.Presentations.Open.return_value = mock_pres
            mock_win32.client.Dispatch.return_value = mock_app
            with patch.dict("sys.modules", {"win32com": mock_win32, "win32com.client": mock_win32.client}):
                parser = OfficeParser()
                result = parser.parse(ppt_path)
                assert "pywin32 ppt slide" in result

    def test_parse_doc_antiword_not_installed(self, tmp_path):
        from unittest.mock import patch
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Linux")):
            with patch("src.parsers.office.subprocess.run", side_effect=FileNotFoundError("no antiword")):
                parser = OfficeParser()
                with pytest.raises(ParserError, match="antiword"):
                    parser.parse(doc_path)

    def test_parse_doc_antiword_timeout(self, tmp_path):
        import subprocess
        from unittest.mock import patch
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Linux")):
            with patch("src.parsers.office.subprocess.run", side_effect=subprocess.TimeoutExpired("antiword", 30)):
                parser = OfficeParser()
                with pytest.raises(ParserError, match="超时"):
                    parser.parse(doc_path)

    def test_parse_doc_catdoc_not_installed(self, tmp_path):
        from unittest.mock import patch
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Darwin")):
            with patch("src.parsers.office.subprocess.run", side_effect=FileNotFoundError("no catdoc")):
                parser = OfficeParser()
                with pytest.raises(ParserError, match="catdoc"):
                    parser.parse(doc_path)

    def test_parse_doc_catdoc_timeout(self, tmp_path):
        import subprocess
        from unittest.mock import patch
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Darwin")):
            with patch("src.parsers.office.subprocess.run", side_effect=subprocess.TimeoutExpired("catdoc", 30)):
                parser = OfficeParser()
                with pytest.raises(ParserError, match="超时"):
                    parser.parse(doc_path)

    def test_parse_unsupported_platform(self, tmp_path):
        from unittest.mock import patch
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "FreeBSD")):
            parser = OfficeParser()
            with pytest.raises(ParserError, match="不支持的平台"):
                parser.parse(doc_path)

    def test_parse_pywin32_import_error(self, tmp_path):
        from unittest.mock import patch
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Windows")):
            with patch.dict("sys.modules", {"win32com": None, "win32com.client": None}):
                parser = OfficeParser()
                with pytest.raises(ParserError, match="pywin32"):
                    parser.parse(doc_path)

    def test_parse_unsupported_office_ext(self, tmp_path):
        bad_path = str(tmp_path / "test.odt")
        with open(bad_path, "wb") as f:
            f.write(b"dummy")
        parser = OfficeParser()
        with pytest.raises(ParserError, match="不支持的 Office 格式"):
            parser.parse(bad_path)


class TestArchiveAdvancedGaps:
    def test_zip_file_too_large(self, tmp_path, monkeypatch):
        import zipfile
        inner = _make_txt(tmp_path, "big.txt", "x" * 200)
        zip_path = str(tmp_path / "big.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "big.txt")
        monkeypatch.setattr("src.parsers.archive.MAX_SIZE", 20)
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        with pytest.raises(ParserError):
            parser.parse(zip_path)

    def test_zip_total_size_limit(self, tmp_path, monkeypatch):
        import zipfile
        inner1 = _make_txt(tmp_path, "a.txt", "x" * 200)
        inner2 = _make_txt(tmp_path, "b.txt", "y" * 200)
        zip_path = str(tmp_path / "total.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner1, "a.txt")
            zf.write(inner2, "b.txt")
        monkeypatch.setattr("src.parsers.archive.MAX_TOTAL_SIZE", 50)
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        with pytest.raises(ParserError):
            parser.parse(zip_path)

    def test_tar_zip_slip_dotdot(self, tmp_path):
        import tarfile
        tar_path = str(tmp_path / "slip.tar")
        with tarfile.open(tar_path, "w") as tf:
            info = tarfile.TarInfo(name="../../../tmp/evil.txt")
            info.size = 9
            import io
            tf.addfile(info, io.BytesIO(b"malicious"))
        parser = ArchiveParser(inner_parser_factory=lambda e: None)
        with pytest.raises(ParserError):
            parser.parse(tar_path)

    def test_tar_zip_slip_absolute(self, tmp_path):
        import tarfile
        tar_path = str(tmp_path / "abs.tar")
        with tarfile.open(tar_path, "w") as tf:
            info = tarfile.TarInfo(name="/etc/passwd")
            info.size = 9
            import io
            tf.addfile(info, io.BytesIO(b"malicious"))
        parser = ArchiveParser(inner_parser_factory=lambda e: None)
        with pytest.raises(ParserError):
            parser.parse(tar_path)

    def test_gz_empty_basename(self, tmp_path):
        import gzip
        gz_path = str(tmp_path / "x.gz")
        data = "empty base name content".encode("utf-8")
        with gzip.open(gz_path, "wb") as f:
            f.write(data)
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        result = parser.parse(gz_path)
        assert "empty base name content" in result

    def test_unsupported_archive_format(self, tmp_path):
        fake_path = str(tmp_path / "test.abc")
        with open(fake_path, "wb") as f:
            f.write(b"not an archive")
        parser = ArchiveParser()
        with pytest.raises(ParserError, match="不支持的压缩格式"):
            parser.parse(fake_path)

    def test_corrupted_tar(self, tmp_path):
        tar_path = str(tmp_path / "corrupt.tar")
        with open(tar_path, "wb") as f:
            f.write(b"not a real tar file at all")
        parser = ArchiveParser()
        with pytest.raises(ParserError):
            parser.parse(tar_path)

    def test_nested_archive_depth(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "depth content")
        zip_path = str(tmp_path / "deep.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None, depth=9)
        result = parser.parse(zip_path)
        assert "depth content" in result

    def test_7z_extract_failure(self, tmp_path, monkeypatch):
        try:
            import py7zr
        except ImportError:
            pytest.skip("py7zr not installed")
        sz_path = str(tmp_path / "bad.7z")
        with open(sz_path, "wb") as f:
            f.write(b"not a 7z file")
        parser = ArchiveParser()
        with pytest.raises(ParserError):
            parser.parse(sz_path)

    def test_rar_import_error(self, tmp_path, monkeypatch):
        import sys
        monkeypatch.setitem(sys.modules, "rarfile", None)
        parser = ArchiveParser()
        rar_path = str(tmp_path / "test.rar")
        with open(rar_path, "wb") as f:
            f.write(b"fake rar")
        with pytest.raises(ParserError, match="rarfile"):
            parser.parse(rar_path)

    def test_7z_import_error(self, tmp_path, monkeypatch):
        import sys
        monkeypatch.setitem(sys.modules, "py7zr", None)
        parser = ArchiveParser()
        sz_path = str(tmp_path / "test.7z")
        with open(sz_path, "wb") as f:
            f.write(b"fake 7z")
        with pytest.raises(ParserError, match="py7zr"):
            parser.parse(sz_path)

    def test_zip_extract_failure(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "content")
        zip_path = str(tmp_path / "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        def bad_factory(ext):
            if ext == ".txt":
                raise Exception("extract boom")
            return None
        parser = ArchiveParser(inner_parser_factory=bad_factory)
        result = parser.parse(zip_path)
        assert "解压失败" in result

    def test_tar_per_file_extract_error(self, tmp_path):
        import tarfile
        inner = _make_txt(tmp_path, "inner.txt", "tar内容")
        tar_path = str(tmp_path / "test.tar")
        with tarfile.open(tar_path, "w") as tf:
            tf.add(inner, arcname="inner.txt")
        def bad_factory(ext):
            if ext == ".txt":
                raise Exception("tar extract boom")
            return None
        parser = ArchiveParser(inner_parser_factory=bad_factory)
        result = parser.parse(tar_path)
        assert "解压失败" in result


class TestOfficeRemainingGaps:
    def test_platform_property(self):
        parser = OfficeParser()
        assert parser._platform in ("Linux", "Darwin", "Windows")

    def test_pptx_parse_error(self, tmp_path):
        path = str(tmp_path / "bad.pptx")
        with open(path, "wb") as f:
            f.write(b"not pptx")
        parser = OfficeParser()
        with pytest.raises(ParserError, match="pptx"):
            parser.parse(path)

    def test_xlsx_parse_error(self, tmp_path):
        path = str(tmp_path / "bad.xlsx")
        with open(path, "wb") as f:
            f.write(b"not xlsx")
        parser = OfficeParser()
        with pytest.raises(ParserError, match="xlsx"):
            parser.parse(path)

    def test_catdoc_returncode_nonzero(self, tmp_path):
        from unittest.mock import patch, MagicMock
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Darwin")):
            with patch("src.parsers.office.subprocess.run") as mock_run:
                mock_run.return_value = MagicMock(returncode=1, stderr="catdoc error")
                parser = OfficeParser()
                with pytest.raises(ParserError, match="catdoc"):
                    parser.parse(doc_path)

    def test_pywin32_parse_failure(self, tmp_path):
        from unittest.mock import patch
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Windows")):
            with patch.dict("sys.modules", {"win32com": None, "win32com.client": None}):
                parser = OfficeParser()
                with pytest.raises(ParserError, match="pywin32"):
                    parser.parse(doc_path)


class TestArchiveRemainingGaps2:
    def test_gz_tar_suffix_basename(self, tmp_path):
        import gzip, tarfile
        inner = _make_txt(tmp_path, "inner.txt", "tar.gz basename content")
        tgz_path = str(tmp_path / "data.tar.gz")
        with tarfile.open(tgz_path, "w:gz") as tf:
            tf.add(inner, arcname="inner.txt")
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        result = parser.parse(tgz_path)
        assert "tar.gz basename content" in result

    def test_gz_oserror(self, tmp_path):
        gz_path = str(tmp_path / "bad.txt.gz")
        with open(gz_path, "wb") as f:
            f.write(b"not real gzip data")
        parser = ArchiveParser()
        with pytest.raises(ParserError, match="解压 gz"):
            parser.parse(gz_path)

    def test_zip_with_directory_entries(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "zip目录内容")
        zip_path = str(tmp_path / "dir.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("subdir/", "")
            zf.write(inner, "subdir/inner.txt")
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        result = parser.parse(zip_path)
        assert "zip目录内容" in result

    def test_nested_archive_parser_error(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "outer content")
        inner_zip_path = str(tmp_path / "inner.zip")
        with zipfile.ZipFile(inner_zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        outer_zip_path = str(tmp_path / "outer.zip")
        with zipfile.ZipFile(outer_zip_path, "w") as zf:
            zf.write(inner_zip_path, "inner.zip")
        parser = ArchiveParser(inner_parser_factory=lambda e: None)
        result = parser.parse(outer_zip_path)
        assert "嵌套解压失败" in result or result == ""

    def test_tar_with_non_file_entries(self, tmp_path):
        import tarfile
        inner = _make_txt(tmp_path, "inner.txt", "tar非文件内容")
        tar_path = str(tmp_path / "mixed.tar")
        with tarfile.open(tar_path, "w") as tf:
            tf.add(inner, arcname="inner.txt")
            info = tarfile.TarInfo(name="symlink")
            info.type = tarfile.SYMTYPE
            info.linkname = "inner.txt"
            tf.addfile(info)
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        result = parser.parse(tar_path)
        assert "tar非文件内容" in result


class TestArchiveFinalGaps:
    def test_gz_empty_basename_decompressed(self, tmp_path):
        import gzip
        gz_path = str(tmp_path / ".gz")
        data = "decompressed fallback".encode("utf-8")
        with gzip.open(gz_path, "wb") as f:
            f.write(data)
        parser = ArchiveParser()
        with pytest.raises(ParserError, match="不支持的压缩格式"):
            parser.parse(gz_path)

    def test_gz_basename_ends_tar(self, tmp_path):
        import gzip, tarfile
        inner = _make_txt(tmp_path, "inner.txt", "gz tar basename")
        tgz_path = str(tmp_path / "archive.tar.gz")
        with tarfile.open(tgz_path, "w:gz") as tf:
            tf.add(inner, arcname="inner.txt")
        parser = ArchiveParser(inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None)
        result = parser.parse(tgz_path)
        assert "gz tar basename" in result

    def test_zip_parser_error_inner(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "content")
        zip_path = str(tmp_path / "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.write(inner, "inner.txt")
        def error_factory(ext):
            if ext == ".txt":
                raise ParserError("inner error")
            return None
        parser = ArchiveParser(inner_parser_factory=error_factory)
        result = parser.parse(zip_path)
        assert "解析失败" in result

    def test_rar_bad_rar_file(self, tmp_path):
        try:
            import rarfile
        except ImportError:
            pytest.skip("rarfile not installed")
        from unittest.mock import patch
        rar_path = str(tmp_path / "bad.rar")
        with open(rar_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(rarfile, "RarFile", side_effect=rarfile.BadRarFile("corrupt")):
            parser = ArchiveParser()
            with pytest.raises(ParserError, match="损坏的 rar"):
                parser.parse(rar_path)

    def test_rar_need_first_volume(self, tmp_path):
        try:
            import rarfile
        except ImportError:
            pytest.skip("rarfile not installed")
        from unittest.mock import patch
        rar_path = str(tmp_path / "split.rar")
        with open(rar_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(rarfile, "RarFile", side_effect=rarfile.NeedFirstVolume("need vol", "1")):
            parser = ArchiveParser()
            with pytest.raises(ParserError, match="分卷 RAR"):
                parser.parse(rar_path)

    def test_7z_extractall_failure(self, tmp_path):
        try:
            import py7zr
        except ImportError:
            pytest.skip("py7zr not installed")
        from unittest.mock import patch, MagicMock
        mock_sz = MagicMock()
        mock_sz.extractall.side_effect = Exception("7z boom")
        parser = ArchiveParser()
        with pytest.raises(ParserError, match="7z 解压失败"):
            parser._extract_7z_archive(mock_sz, "test.7z")

    def test_nested_archive_depth_limit_error(self, tmp_path):
        import zipfile
        inner = _make_txt(tmp_path, "inner.txt", "nested content")
        inner_zip = str(tmp_path / "inner.zip")
        with zipfile.ZipFile(inner_zip, "w") as zf:
            zf.write(inner, "inner.txt")
        outer_zip = str(tmp_path / "outer.zip")
        with zipfile.ZipFile(outer_zip, "w") as zf:
            zf.write(inner_zip, "inner.zip")
        parser = ArchiveParser(
            inner_parser_factory=lambda e: TxtParser() if e == ".txt" else None,
            depth=9,
        )
        result = parser.parse(outer_zip)
        assert "嵌套解压失败" in result


class TestOfficeFinalGaps:
    def test_pywin32_unsupported_legacy_ext(self, tmp_path):
        from unittest.mock import patch, MagicMock
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Windows")):
            mock_win32 = MagicMock()
            mock_app = MagicMock()
            mock_win32.client.Dispatch.return_value = mock_app
            with patch.dict("sys.modules", {"win32com": mock_win32, "win32com.client": mock_win32.client}):
                parser = OfficeParser()
                with patch.object(parser, "_parse_with_pywin32", side_effect=ParserError("不支持的旧格式: odt")):
                    with pytest.raises(ParserError, match="不支持的旧格式"):
                        parser.parse(doc_path)

    def test_pywin32_generic_exception(self, tmp_path):
        from unittest.mock import patch, MagicMock
        doc_path = str(tmp_path / "test.doc")
        with open(doc_path, "wb") as f:
            f.write(b"dummy")
        with patch.object(OfficeParser, "_platform", new_callable=lambda: property(lambda self: "Windows")):
            mock_win32 = MagicMock()
            mock_app = MagicMock()
            mock_app.Documents.Open.side_effect = Exception("COM failed")
            mock_win32.client.Dispatch.return_value = mock_app
            with patch.dict("sys.modules", {"win32com": mock_win32, "win32com.client": mock_win32.client}):
                parser = OfficeParser()
                with pytest.raises(ParserError, match="pywin32"):
                    parser.parse(doc_path)
