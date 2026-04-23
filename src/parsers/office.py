import os
import platform
import subprocess
import sys

from .base import BaseParser, ParserError


class OfficeParser(BaseParser):
    _PLATFORM = platform.system()

    def parse(self, file_path: str) -> str:
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        if ext == ".docx":
            return self._parse_docx(file_path)
        elif ext == ".pptx":
            return self._parse_pptx(file_path)
        elif ext == ".xlsx":
            return self._parse_xlsx(file_path)
        elif ext == ".doc":
            return self._parse_doc(file_path)
        elif ext == ".xls":
            return self._parse_xls(file_path)
        elif ext == ".ppt":
            return self._parse_ppt(file_path)
        else:
            raise ParserError(f"不支持的 Office 格式: {ext}")

    def _parse_docx(self, file_path: str) -> str:
        try:
            from docx import Document

            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            raise ParserError(f"解析 docx 失败: {file_path}: {e}") from e

    def _parse_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            texts.append(para.text)
            return "\n".join(texts)
        except Exception as e:
            raise ParserError(f"解析 pptx 失败: {file_path}: {e}") from e

    def _parse_xlsx(self, file_path: str) -> str:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            texts.append(str(cell))
            wb.close()
            return "\n".join(texts)
        except Exception as e:
            raise ParserError(f"解析 xlsx 失败: {file_path}: {e}") from e

    def _parse_doc(self, file_path: str) -> str:
        return self._parse_legacy(file_path, "doc")

    def _parse_xls(self, file_path: str) -> str:
        return self._parse_legacy(file_path, "xls")

    def _parse_ppt(self, file_path: str) -> str:
        return self._parse_legacy(file_path, "ppt")

    def _parse_legacy(self, file_path: str, ext: str) -> str:
        if self._PLATFORM == "Linux":
            return self._parse_with_antiword(file_path)
        elif self._PLATFORM == "Darwin":
            return self._parse_with_catdoc(file_path)
        elif self._PLATFORM == "Windows":
            return self._parse_with_pywin32(file_path, ext)
        else:
            raise ParserError(f"不支持的平台: {self._PLATFORM}")

    def _parse_with_antiword(self, file_path: str) -> str:
        try:
            result = subprocess.run(
                ["antiword", file_path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode != 0:
                raise ParserError(f"antiword 解析失败: {result.stderr.strip()}")
            return result.stdout
        except FileNotFoundError:
            raise ParserError("antiword 未安装，请执行: apt install antiword")
        except subprocess.TimeoutExpired:
            raise ParserError("antiword 解析超时")

    def _parse_with_catdoc(self, file_path: str) -> str:
        try:
            result = subprocess.run(
                ["catdoc", file_path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode != 0:
                raise ParserError(f"catdoc 解析失败: {result.stderr.strip()}")
            return result.stdout
        except FileNotFoundError:
            raise ParserError("catdoc 未安装，请执行: brew install catdoc")
        except subprocess.TimeoutExpired:
            raise ParserError("catdoc 解析超时")

    def _parse_with_pywin32(self, file_path: str, ext: str) -> str:
        try:
            import win32com.client

            app_map = {
                "doc": "Word.Application",
                "xls": "Excel.Application",
                "ppt": "PowerPoint.Application",
            }
            app_name = app_map.get(ext)
            if not app_name:
                raise ParserError(f"不支持的旧格式: {ext}")

            app = win32com.client.Dispatch(app_name)
            app.Visible = False
            try:
                abs_path = os.path.abspath(file_path)
                if ext == "doc":
                    doc = app.Documents.Open(abs_path)
                    text = doc.Content.Text
                    doc.Close(False)
                elif ext == "xls":
                    wb = app.Workbooks.Open(abs_path)
                    texts = []
                    for ws in wb.Worksheets:
                        used = ws.UsedRange
                        if used:
                            for row in used:
                                for cell in row:
                                    if cell.Value is not None:
                                        texts.append(str(cell.Value))
                    text = "\n".join(texts)
                    wb.Close(False)
                elif ext == "ppt":
                    pres = app.Presentations.Open(abs_path)
                    texts = []
                    for slide in pres.Slides:
                        for shape in slide.Shapes:
                            if shape.HasTextFrame:
                                texts.append(shape.TextFrame.TextRange.Text)
                    text = "\n".join(texts)
                    pres.Close()
                else:
                    text = ""
                return text
            finally:
                app.Quit()
        except ImportError:
            raise ParserError("pywin32 未安装，请执行: pip install pywin32")
        except Exception as e:
            raise ParserError(f"pywin32 解析失败: {e}")
